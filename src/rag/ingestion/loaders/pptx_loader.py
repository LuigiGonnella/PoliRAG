"""PowerPoint slide structural text parsing and window chunking routine."""
from pptx import Presentation

def chunk_pptx(file_path, chunk_size=500, chunk_overlap=50):
    """
    Extracts text shapes, titles, and layout containers slide-by-slide, 
    assembling uniform rolling token chunk arrays.
    """
    prs = Presentation(file_path)
    slide_texts = []

    # 1. Gather all slide structural contents text fields
    for slide_idx, slide in enumerate(prs.slides):
        current_slide_content = []
        
        # Pull text from all layout boxes, text frames, or graphics shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                current_slide_content.append(shape.text.strip())
                
        # Optional: Extract speaker presentation notes if they exist
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                current_slide_content.append(f"[Notes: {notes}]")
                
        if current_slide_content:
            combined_slide_text = f" [Slide {slide_idx + 1}] " + " | ".join(current_slide_content)
            slide_texts.append(combined_slide_text)

    full_document_text = "\n".join(slide_texts)
    words = full_document_text.split()
    
    chunks = []
    chunk_index = 0
    
    # 2. Build rolling sliding window chunks matching system schemas
    i = 0
    while i < len(words):
        window_words = words[i:i + chunk_size]
        if not window_words:
            break
            
        chunk_text = " ".join(window_words)
        chunks.append({
            "text": chunk_text,
            "source": file_path,
            "index": chunk_index
        })
        
        chunk_index += 1
        i += (chunk_size - chunk_overlap)
        
    return chunks